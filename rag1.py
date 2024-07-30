# from typing import List
# import PyPDF2
# from io import BytesIO
# from langchain_community.embeddings import OllamaEmbeddings
# from langchain_openai import OpenAIEmbeddings
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_community.vectorstores.faiss import FAISS
# from langchain.chains import ConversationalRetrievalChain
# from langchain.docstore.document import Document
# from langchain_community.chat_models import ChatOllama
# from langchain.memory import ConversationBufferMemory
# from langchain_community.chat_message_histories import ChatMessageHistory
# import chainlit as cl
# import docx
# import pptx

# text_splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=100)

# @cl.on_chat_start
# async def on_chat_start():
#     all_file_texts = []  # Initialize a list to hold text from all files

#     # Wait for the user to upload multiple files
#     uploaded_files = await cl.AskFileMessage(
#         content="Please upload pdf, docx, txt, or pptx files! You can select multiple files.",
#         accept=["application/pdf", 
#                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
#                 "text/plain", 
#                 "application/vnd.openxmlformats-officedocument.presentationml.presentation"],
#         max_size_mb=20,
#         max_files=10  # Allow multiple file uploads
#     ).send()

#     for uploaded_file in uploaded_files:
#         print(f"Uploaded file: {uploaded_file.name}")

#         msg = cl.Message(content=f"Processing `{uploaded_file.name}`...")
#         await msg.send()

#         # Read the uploaded file based on its type
#         file_ext = uploaded_file.name.split('.')[-1].lower()
#         file_text = ""

#         if file_ext == 'pdf':
#             print("Reading PDF file...")
#             pdf = PyPDF2.PdfReader(uploaded_file.path)
#             for page in pdf.pages:
#                 file_text += page.extract_text() or ""  # Handle cases where extract_text returns None
#         elif file_ext == 'docx':
#             print("Reading DOCX file...")
#             doc = docx.Document(uploaded_file.path)
#             for para in doc.paragraphs:
#                 file_text += para.text + '\n'
#         elif file_ext == 'txt':
#             print("Reading TXT file...")
#             with open(uploaded_file.path, 'r', encoding='utf-8') as f:
#                 file_text = f.read()
#         elif file_ext == 'pptx':
#             print("Reading PPTX file...")
#             prs = pptx.Presentation(uploaded_file.path)
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         file_text += shape.text + '\n'
#         else:
#             print(f"Unsupported file type: {file_ext}")
#             continue  # Skip unsupported file types

#         print(f"Document text extracted from `{uploaded_file.name}`. Total characters: {len(file_text)}")
#         all_file_texts.append(file_text)  # Append the text to the list

#     # Combine text from all files
#     combined_text = "\n".join(all_file_texts)

#     # Split the combined text into chunks
#     print("Splitting text into chunks...")
#     texts = text_splitter.split_text(combined_text)
#     print(f"Number of text chunks: {len(texts)}")

#     # Create metadata for each chunk
#     metadatas = [{"source": f"Combined-File-{i}"} for i in range(len(texts))]

#     print("Creating FAISS vector store...")
#     embeddings = OllamaEmbeddings(model="gemma:2b")
#     docsearch = FAISS.from_texts(
#         texts, embeddings, metadatas=metadatas
#     )
#     print("FAISS vector store created successfully.")

#     message_history = ChatMessageHistory()
#     memory = ConversationBufferMemory(
#         memory_key="chat_history",
#         output_key="answer",
#         chat_memory=message_history,
#         return_messages=True,
#     )

#     # Create a chain that uses the FAISS vector store
#     print("Setting up conversational chain...")
#     chain = ConversationalRetrievalChain.from_llm(
#         ChatOllama(model="mistral"),
#         chain_type="stuff",
#         retriever=docsearch.as_retriever(),
#         memory=memory,
#         return_source_documents=True,
#     )

#     # Let the user know that the system is ready
#     print(f"Processing done. You can now ask questions!")
#     msg.content = "Processing done. You can now ask questions!"
#     await msg.update()

#     cl.user_session.set("chain", chain)

# @cl.on_message
# async def main(message: cl.Message):
#     chain = cl.user_session.get("chain")  # type: ConversationalRetrievalChain

#     cb = cl.AsyncLangchainCallbackHandler()
#     res = await chain.ainvoke(message.content, callbacks=[cb])
#     answer = res["answer"]
#     source_documents = res["source_documents"]  # type: List[Document]

#     text_elements = []  # type: List[cl.Text]

#     if source_documents:
#         for source_idx, source_doc in enumerate(source_documents):
#             source_name = f"source_{source_idx}"
#             # Create the text element referenced in the message
#             text_elements.append(
#                 cl.Text(content=source_doc.page_content, name=source_name)
#             )
#         source_names = [text_el.name for text_el in text_elements]

#         if source_names:
#             answer += f"\nSources: {', '.join(source_names)}"
#         else:
#             answer += "\nNo sources found"

#     await cl.Message(content=answer, elements=text_elements).send()

from typing import List
import PyPDF2
from io import BytesIO
from langchain_community.embeddings import OllamaEmbeddings
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.docstore.document import Document
from langchain_community.chat_models import ChatOllama
from langchain.memory import ConversationBufferMemory
from langchain_community.chat_message_histories import ChatMessageHistory
import chainlit as cl
import docx
import pptx

text_splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=100)

@cl.on_chat_start
async def on_chat_start():
    all_file_texts = []  # Initialize a list to hold text from all files

    # Wait for the user to upload multiple files
    uploaded_files = await cl.AskFileMessage(
        content="Please upload pdf, docx, txt, or pptx files! You can select multiple files.",
        accept=["application/pdf", 
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                "text/plain", 
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"],
        max_size_mb=20,
        max_files=10  # Allow multiple file uploads
    ).send()

    for uploaded_file in uploaded_files:
        print(f"Uploaded file: {uploaded_file.name}")

        msg = cl.Message(content=f"Processing `{uploaded_file.name}`...")
        await msg.send()

        # Read the uploaded file based on its type
        file_ext = uploaded_file.name.split('.')[-1].lower()
        file_text = ""

        if file_ext == 'pdf':
            print("Reading PDF file...")
            pdf = PyPDF2.PdfReader(uploaded_file.path)
            for page in pdf.pages:
                file_text += page.extract_text() or ""  # Handle cases where extract_text returns None
            
            # Display the PDF using the Pdf class
            elements = [
                cl.Pdf(name=uploaded_file.name, display="side", path=uploaded_file.path)
            ]
            await cl.Message(content=f"Displaying your PDF: {uploaded_file.name}", elements=elements).send()

        elif file_ext == 'docx':
            print("Reading DOCX file...")
            doc = docx.Document(uploaded_file.path)
            for para in doc.paragraphs:
                file_text += para.text + '\n'
        elif file_ext == 'txt':
            print("Reading TXT file...")
            with open(uploaded_file.path, 'r', encoding='utf-8') as f:
                file_text = f.read()
        elif file_ext == 'pptx':
            print("Reading PPTX file...")
            prs = pptx.Presentation(uploaded_file.path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        file_text += shape.text + '\n'
        else:
            print(f"Unsupported file type: {file_ext}")
            continue  # Skip unsupported file types

        print(f"Document text extracted from `{uploaded_file.name}`. Total characters: {len(file_text)}")
        all_file_texts.append(file_text)  # Append the text to the list

    # Combine text from all files
    combined_text = "\n".join(all_file_texts)

    # Split the combined text into chunks
    print("Splitting text into chunks...")
    texts = text_splitter.split_text(combined_text)
    print(f"Number of text chunks: {len(texts)}")

    # Create metadata for each chunk
    metadatas = [{"source": f"Combined-File-{i}"} for i in range(len(texts))]

    print("Creating FAISS vector store...")
    embeddings = OllamaEmbeddings(model="gemma:2b")
    docsearch = FAISS.from_texts(
        texts, embeddings, metadatas=metadatas
    )
    print("FAISS vector store created successfully.")

    message_history = ChatMessageHistory()
    memory = ConversationBufferMemory(
        memory_key="chat_history",
        output_key="answer",
        chat_memory=message_history,
        return_messages=True,
    )

    # Create a chain that uses the FAISS vector store
    print("Setting up conversational chain...")
    chain = ConversationalRetrievalChain.from_llm(
        ChatOllama(model="mistral"),
        chain_type="stuff",
        retriever=docsearch.as_retriever(),
        memory=memory,
        return_source_documents=True,
    )

    # Let the user know that the system is ready
    print(f"Processing done. You can now ask questions!")
    msg.content = "Processing done. You can now ask questions!"
    await msg.update()

    cl.user_session.set("chain", chain)

    
@cl.on_message
async def main(message: cl.Message):
    chain = cl.user_session.get("chain")  # type: ConversationalRetrievalChain

    cb = cl.AsyncLangchainCallbackHandler()
    res = await chain.ainvoke(message.content, callbacks=[cb])
    answer = res["answer"]
    source_documents = res["source_documents"]  # type: List[Document]

    text_elements = []  # type: List[cl.Text]

    if source_documents:
        for source_idx, source_doc in enumerate(source_documents):
            source_name = f"source_{source_idx}"
            # Create the text element referenced in the message
            text_elements.append(
                cl.Text(content=source_doc.page_content, name=source_name)
            )
        source_names = [text_el.name for text_el in text_elements]

        if source_names:
            answer += f"\nSources: {', '.join(source_names)}"
        else:
            answer += "\nNo sources found"

    await cl.Message(content=answer, elements=text_elements).send()

